"""
多格式文件解析器 — 支持 PDF/Word/PPT/Excel/TXT/MD → 纯文本
"""
from pathlib import Path


def parse_pdf(file_path: Path) -> str:
    """提取 PDF 文本（使用 pymupdf）"""
    import fitz  # pymupdf
    doc = fitz.open(str(file_path))
    try:
        texts = []
        for page in doc:
            t = page.get_text("text")
            if t:
                texts.append(t)
        return "\n".join(texts)
    finally:
        doc.close()


def parse_docx(file_path: Path) -> str:
    """提取 Word .docx 文本（使用 python-docx）"""
    from docx import Document
    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def parse_pptx(file_path: Path) -> str:
    """提取 PowerPoint .pptx 文本（使用 python-pptx）"""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return "\n".join(texts)


def parse_xlsx(file_path: Path) -> str:
    """提取 Excel .xlsx 文本（使用 openpyxl）"""
    import openpyxl
    wb = openpyxl.load_workbook(str(file_path), data_only=True)
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        texts.append(f"Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                texts.append(row_text)
    wb.close()
    return "\n".join(texts)


def parse_txt(file_path: Path) -> str:
    """读取纯文本 TXT 文件（UTF-8，失败回退 GBK）"""
    try:
        return file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return file_path.read_text(encoding="gbk")


def parse_md(file_path: Path) -> str:
    """读取 Markdown 文件（同 TXT）"""
    return parse_txt(file_path)


# 扩展名 → 解析函数映射
_PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_xlsx,
    ".txt": parse_txt,
    ".md": parse_md,
}


def parse_file(file_path: Path, file_type: str) -> str:
    """统一入口：根据扩展名调用对应解析器

    Args:
        file_path: 文件路径
        file_type: 小写扩展名（含点），如 ".pdf"

    Returns:
        提取的纯文本字符串；空文件或无文本内容返回空字符串

    Raises:
        ValueError: 不支持的文件类型
    """
    parser = _PARSERS.get(file_type.lower())
    if parser is None:
        raise ValueError(f"不支持的文件类型: {file_type}")
    return parser(file_path)
